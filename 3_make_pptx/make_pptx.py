import os
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../utils/'))) # allow importing the utils dir
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../2_graph_processing/'))) # allow importing the graphs dir

import argparse

import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.subplots as ps


import common_utils
import constants
import make_pptx_helpers
import make_pptx_legend

import pptx
import pptx.util
import pptx.enum.shapes
import pptx.enum.text
import pptx.dml.color
import pptx.dml.effect
import PIL

PPTX_TEMPLATE_FILE = os.path.join(os.path.dirname(__file__), 'template.pptx') # make this an arg!

FORMAT_LEGENDS = {
  'combined': [
    'node_size',
    'node_outline',
    'edge_type',
  ],
  'individual': [
    'node_size',
    'node_outline',
    'edge_type',
    'variation_type',
  ],
  'both': [
    'node_size',
    'node_outline',
    'edge_type',
    'variation_type',
  ],
}

FREQ_RATIO_LEGENDS = {
  '1DSB': [
    'freq_ratio_sense_branch',
    'freq_ratio_sense_cmv',
  ],
  '2DSB': [
    'freq_ratio_sense_branch',
    'freq_ratio_sense_cmv',
  ],
  '2DSBanti': [
    'freq_ratio_antisense_splicing',
  ],
}

TITLE_FONT_SIZE_PT = 16
MARGIN_CORNER_LABEL_FONT_SIZE_PT = 12
MARGIN_TOP_LABEL_FONT_SIZE_PT = 9
MARGIN_LEFT_LABEL_FONT_SIZE_PT = 8
LEGEND_TITLE_FONT_SIZE_PT = 10
LEGEND_LABEL_FONT_SIZE_PT = 8

TITLE_HEIGHT_PT = 30
MARGIN_TOP_HEIGHT_PT = 20
MARGIN_LEFT_WIDTH_PT = 60
MARGIN_LEFT_SPILL_OVER_PT = 10
MARGIN_RIGHT_WIDTH_PT = 100

LEGEND_TITLE_HEIGHT_PT = 20
LEGEND_ITEM_HEIGHT_PT = 20
LEGEND_SIZE_NODE_OUTLINE_WIDTH_PT = 0.1
LEGEND_NODE_SIZE_PT = 10
LEGEND_NODE_OUTLINE_WIDTH_PT = 0.5
LEGEND_EDGE_LINE_SIZE_PT = 10
LEGEND_EDGE_LINE_WIDTH_PT = 1
LEGEND_FREQ_RATIO_COLOR_BAR_WIDTH_PT = 10
LEGEND_FREQ_RATIO_COLOR_BAR_HEIGHT_PT = 75

LEGEND_HEIGHT_SPACING_PT = 25

HEIGHT_SPACING_PT = 10
WIDTH_SPACING_PT = 5
MULTIPLE_GRIDS_SPACING_PT = 40
CONTENT_LEGEND_SPACING_PT = 20


def get_data_set_spec_label(spec):
  labels_list = []

  if spec['DSB'] == '2DSB':
    labels_list.append(common.LABELS['sgAB'])
  elif spec['DSB'] == '2DSBanti':
    labels_list.append(common.LABELS['sgCD'])
  
  if spec['DSB'] in ['2DSB', '2DSBanti']:
    labels_list.append(common.LABELS[spec['strand']])
  elif spec['DSB'] in ['1DSB']:
    labels_list.append(common.LABELS[spec['hguide']])
  else:
    raise Exception('Bad DSB key: ' + str(spec['DSB']))
  
  if spec['control'] != 'not_control':
    labels_list.append(common.LABELS[spec['control']])
  
  return join_margin_label(labels_list)

def make_slide(
  prs,
  title,
  image_grid_list,
  label_grid_list,
  node_size_min_freq = None,
  node_size_max_freq = None,
  node_size_min_px = None,
  node_size_max_px = None,
  legend_list = [],
  legend_height_spacing_pt = LEGEND_HEIGHT_SPACING_PT,
  legend_freq_ratio_color_bar_height_pt = LEGEND_FREQ_RATIO_COLOR_BAR_HEIGHT_PT,
  legend_title_font_size_pt = LEGEND_TITLE_FONT_SIZE_PT,
  legend_label_font_size_pt = LEGEND_LABEL_FONT_SIZE_PT,
):
  title_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(title_slide_layout)
  
  slide_width_pt = prs.slide_width / pptx.util.Pt(1)

  image_label_width_pt = 60
  image_label_height_pt = 20
  
  y_pt = 0
  x_pt = 0
  y_legend_pt = 0
  x_legend_pt = slide_width_pt

  ## Legend constants ###
  legend_title_width_pt = 100
  legend_title_height_pt = LEGEND_TITLE_HEIGHT_PT
  legend_item_width_pt = 30
  legend_item_height_pt = LEGEND_ITEM_HEIGHT_PT
  legend_label_width_pt = 70
  legend_label_height_pt = LEGEND_ITEM_HEIGHT_PT

  legend_const = {
    'v': {
      'node_size': {
        'stride_pt': 20,
        'title_width_pt': 100,
        'title_height_pt': 20,
        'item_width_pt': 30,
        'label_width_pt': 50,
        'item_height_pt': 20,
        'label_height_pt': 20,
      },
      'node_outline': {
        'stride_pt': 20,
        'title_width_pt': 100,
        'title_height_pt': 20,
        'item_width_pt': 30,
        'label_width_pt': 50,
        'item_height_pt': 20,
        'label_height_pt': 20,
      },
      'edge_type': {
        'stride_pt': 20,
        'title_width_pt': 100,
        'title_height_pt': 20,
        'item_width_pt': 30,
        'label_width_pt': 50,
        'item_height_pt': 20,
        'label_height_pt': 20,
      },
    },
    'h': {
      'node_size': {
        'stride_pt': 20,
        'title_width_pt': 100,
        'title_height_pt': 20,
        'item_width_pt': 30,
        'label_width_pt': 50,
        'item_height_pt': 20,
        'label_height_pt': 20,
      }
    },
  }
  legend_x_offset_pt = {
    'v': 0,
    'h': 100,
  }

  for i in range(len(image_grid_list)):
    max_image_width_px = 0
    max_image_height_px = 0

    # Get point/pixel ratio and dimensions
    # All images should be the same width/height
    for row in range(image_grid_list[i].shape[0]):
      for col in range(image_grid_list[i].shape[1]):
        image = PIL.Image.open(image_grid_list[i][row][col])
        max_image_width_px = max(max_image_width_px, image.width)
        max_image_height_px = max(max_image_height_px, image.height)

    content_width_px = image_grid_list[i].shape[1] * max_image_width_px

    ratio_pt_px = slide_width_pt / content_width_px
    cell_width_pt = ratio_pt_px * max_image_width_px
    cell_height_pt = ratio_pt_px * max_image_height_px

    # Content images
    make_pptx_helpers.add_picture_grid_pptx(
      slide = slide,
      file_name_grid = image_grid_list[i],
      x_pt = x_pt,
      y_pt = y_pt,
      cell_width_pt = cell_width_pt,
      cell_height_pt = cell_height_pt,
      cell_width_spacing_pt = 0,
      cell_height_spacing_pt = 0,
    )

    # Content labels
    make_pptx_helpers.add_text_grid_pptx(
      slide = slide,
      text_grid = label_grid_list[i],
      x_pt = x_pt,
      y_pt = y_pt,
      cell_width_pt = cell_width_pt,
      cell_height_pt = cell_height_pt,
      cell_width_spacing_pt = 0,
      cell_height_spacing_pt = 0,
      text_width_pt = image_label_width_pt,
      text_height_pt = image_label_height_pt,
      font_size_pt = MARGIN_TOP_LABEL_FONT_SIZE_PT,
      text_align = 'center',
    )
    # for row in range(image_grid_list[i].shape[0]):
    #   for col in range(image_grid_list[i].shape[1]):
    #     make_pptx_helpers.add_textbox_pptx(
    #       slide = slide,
    #       text = label_grid_list[i][row, col],
    #       x_pt = x_pt + (col) * cell_width_pt,
    #       y_pt = y_pt + (row) * cell_height_pt,
    #       width_pt = image_label_width_pt,
    #       height_pt = image_label_height_pt,
    #       font_size_pt = MARGIN_TOP_LABEL_FONT_SIZE_PT,
    #     )

    y_pt += image_grid_list[i].shape[0] * cell_height_pt
    
    # Size legends
    for orientation in ['h', 'v']:
      y_legend_new_pt = make_pptx_legend.make_size_legend_pptx(
        slide = slide,
        x_pt = x_legend_pt + legend_x_offset_pt[orientation],
        y_pt = y_legend_pt,
        stride_pt = legend_const[orientation]['node_size']['stride_pt'],
        title_width_pt = legend_const[orientation]['node_size']['title_width_pt'],
        title_height_pt = legend_const[orientation]['node_size']['title_height_pt'],
        item_width_pt = legend_const[orientation]['node_size']['item_width_pt'],
        item_height_pt = legend_const[orientation]['node_size']['item_height_pt'],
        label_width_pt = legend_const[orientation]['node_size']['label_width_pt'],
        label_height_pt = legend_const[orientation]['node_size']['label_height_pt'],
        node_size_min_freq = node_size_min_freq,
        node_size_max_freq = node_size_max_freq,
        node_size_min_pt = node_size_min_px * ratio_pt_px,
        node_size_max_pt = node_size_max_px * ratio_pt_px,
        line_width_pt = LEGEND_SIZE_NODE_OUTLINE_WIDTH_PT,
        legend_title_font_size_pt = legend_title_font_size_pt,
        legend_label_font_size_pt = legend_label_font_size_pt,
        orientation = orientation,
      )
    y_legend_pt = y_legend_new_pt

  x_pt = slide_width_pt # place legends outside the actual slide

  # Legends
  for legend in legend_list:
    for orientation in ['h', 'v']:
      if legend['type'] == 'node_size':
        pass # already handled above
      elif legend['type'] == 'node_outline':
        y_legend_new_pt = make_pptx_legend.make_outline_legend_pptx(
          slide = slide,
          x_pt = x_pt + legend_x_offset_pt[orientation],
          y_pt = y_legend_pt,
          stride_pt = legend_item_height_pt,
          title_width_pt = legend_title_width_pt,
          title_height_pt = legend_title_height_pt,
          item_width_pt = legend_item_width_pt,
          item_height_pt = legend_item_height_pt,
          label_width_pt = legend_label_width_pt,
          label_height_pt = legend_item_height_pt,
          node_size_pt = LEGEND_NODE_SIZE_PT,
          line_width_pt = LEGEND_NODE_OUTLINE_WIDTH_PT,
          legend_title_font_size_pt = legend_title_font_size_pt,
          legend_label_font_size_pt = legend_label_font_size_pt,
          orientation = orientation,
        )
      elif legend['type'] == 'edge_type':
        y_legend_new_pt = make_pptx_legend.make_edge_legend_pptx(
          slide = slide,
          x_pt = x_pt + legend_x_offset_pt[orientation],
          y_pt = y_legend_pt,
          stride_pt = legend_item_height_pt,
          title_width_pt = legend_title_width_pt,
          title_height_pt = legend_title_height_pt,
          item_width_pt = legend_item_width_pt,
          item_height_pt = legend_item_height_pt,
          label_width_pt = legend_label_width_pt,
          label_height_pt = legend_item_height_pt,
          line_size_pt = LEGEND_EDGE_LINE_SIZE_PT,
          line_width_pt = LEGEND_EDGE_LINE_WIDTH_PT,
          legend_title_font_size_pt = legend_title_font_size_pt,
          legend_label_font_size_pt = legend_label_font_size_pt,
          orientation = orientation,
        )
      elif legend['type'] == 'variation_type':
        y_legend_new_pt = make_pptx_legend.make_variation_color_legend_pptx(
          slide = slide,
          x_pt = x_pt + legend_x_offset_pt[orientation],
          y_pt = y_legend_pt,
          stride_pt = legend_item_height_pt,
          title_width_pt = legend_title_width_pt,
          title_height_pt = legend_title_height_pt,
          item_width_pt = legend_item_width_pt,
          item_height_pt = legend_item_height_pt,
          label_width_pt = legend_label_width_pt,
          label_height_pt = legend_label_height_pt,
          variation_types = ['insertion', 'deletion', 'none'],
          node_size_pt = LEGEND_NODE_SIZE_PT,
          line_width_pt = LEGEND_NODE_OUTLINE_WIDTH_PT,
          legend_title_font_size_pt = legend_title_font_size_pt,
          legend_label_font_size_pt = legend_label_font_size_pt,
          orientation = orientation,
        )
      elif legend['type'] == 'freq_ratio':
        treatment_1 = legend['treatment_1']
        treatment_2 = legend['treatment_2']
        title = f'Ratio [{constants.LABELS[treatment_1]} / {constants.LABELS[treatment_2]}]'
        color_bar_file = legend['color_bar_file']
        y_legend_new_pt = make_pptx_legend.make_freq_ratio_legend_pptx(
          slide = slide,
          x_pt = x_pt + legend_x_offset_pt[orientation],
          y_pt = y_legend_pt,
          title = title,
          title_width_pt = legend_title_width_pt,
          title_height_pt = legend_title_height_pt * 2,
          label_width_pt = legend_label_width_pt,
          label_height_pt = legend_label_height_pt,
          label_width_left_spill_over_pt = 10,
          color_bar_minor_axis_pt = LEGEND_FREQ_RATIO_COLOR_BAR_WIDTH_PT,
          color_bar_major_axis_pt = legend_freq_ratio_color_bar_height_pt,
          color_bar_file = color_bar_file,
          title_pad_pt = 7 if treatment_1 == 'antisense_splicing' else 5,
          legend_title_font_size_pt = legend_title_font_size_pt,
          legend_label_font_size_pt = legend_label_font_size_pt,
          orientation = orientation,
        )
      else:
        raise Exception('Unknown legend type: ' + str(legend))
    y_legend_pt = y_legend_new_pt
    y_legend_pt += legend_height_spacing_pt


def parse_args():
  parser = argparse.ArgumentParser(
    description = 'Create powerpoint figures from the graphs.'
  )
  parser.add_argument(
    '-i',
    '--input',
    nargs = '+',
    type = argparse.FileType('r'),
    help = 'List of images to include in the grid',
    required = True,
  )
  parser.add_argument(
    '-o',
    '--output',
    type = argparse.FileType('w'),
    help = 'Output PPTX file.',
    required = True,
  )
  parser.add_argument(
    '-ng',
    '--num_grids',
    required = True,
    type = int,
    help = 'Number of separate grids to create',
  )
  parser.add_argument(
    '-nr',
    '--num_rows',
    nargs = '+',
    required = True,
    type = int,
    help = (
      'Number of rows in each grid.' +
      ' Number of arguments should match the number of grids.'
    ),
  )
  parser.add_argument(
    '-nc',
    '--num_cols',
    nargs = '+',
    required = True,
    type = int,
    help = (
      'Number of columns in each grid.' +
      ' Number of arguments should match the number of grids.'
    ),
  )
  # parser.add_argument(
  #   '--node_max_freq',
  #   type = float,
  #   help = (
  #     'Max frequency to determine node size.' +
  #     'Higher frequencies are clipped to this value.'
  #   ),
  #   default = constants.GRAPH_NODE_SIZE_MAX_FREQ,
  # )
  # parser.add_argument(
  #   '--node_min_freq',
  #   type = float,
  #   help = (
  #     'Min frequency to determine node size.' +
  #     'Lower frequencies are clipped to this value.'
  #   ),
  #   default = constants.GRAPH_NODE_SIZE_MIN_FREQ,
  # )
  # parser.add_argument(
  #   '--node_max_px',
  #   type = float,
  #   help = 'Largest node size as determined by the frequency.',
  #   default = constants.GRAPH_NODE_SIZE_MAX_PX,
  # )
  # parser.add_argument(
  #   '--node_min_px',
  #   type = float,
  #   help = 'Smallest node size as determined by the frequency.',
  #   default = constants.GRAPH_NODE_SIZE_MIN_PX,
  # )
  
  return parser.parse_args()


if __name__ == '__main__':
  # make_legend_images()
  # make_pptx_1()
  # args = parse_args()

  # if args.num_grids != len(args.num_rows):
  #   raise Exception(
  #     f'Incorrect num rows specification: {args.num_rows}.' +
  #     f' Expected {args.num_grids} values.' 
  #   )
  
  # if args.num_grids != len(args.num_columns):
  #   raise Exception(
  #     f'Incorrect num columns specification: {args.num_cols}.' +
  #     f' Expected {args.num_grids} values.' 
  #   )

  # num_images_total = sum(r * c for r, c in zip(args.num_rows, args.num_cols))
  # if num_images_total != len(args.input):
  #   raise Exception(
  #     f'Incorrect number of input files: {len(args.input)}.' +
  #     f' Expected {num_images_total} values.' 
  #   )

  # image_grid_list = []
  # image_index = 0
  # for i in range(args.num_grids):
  #   num_images = args.num_rows[i] * args.num_cols[i]
  #   image_grid = np.array(
  #     [args.input[image_index + j] for j in range(num_images)]
  #   )
  #   image_grid = image_grid.reshape((args.num_rows[i], args.num_cols[i]))
  #   image_grid_list.append(image_grid)
  #   image_index += num_images

  
  prs = pptx.Presentation(PPTX_TEMPLATE_FILE)
  image_grid = np.array(
    [
      "WT_sgAB_R1_branch.png", "WT_sgAB_R1_cmv.png", "WT_sgAB_R1_sense.png",
      "WT_sgAB_R2_branch.png", "WT_sgAB_R2_cmv.png", "WT_sgAB_R2_sense.png",
      "WT_sgA_R1_branch.png", "WT_sgA_R1_cmv.png", "WT_sgA_R1_sense.png",
      "WT_sgB_R2_branch.png", "WT_sgB_R2_cmv.png", "WT_sgB_R2_sense.png",
    ],
    dtype = object,
  ).reshape((-1, 3))
  label_grid = np.array(
    [
      "WT_sgAB_R1_branch", "WT_sgAB_R1_cmv", "WT_sgAB_R1_sense",
      "WT_sgAB_R2_branch", "WT_sgAB_R2_cmv", "WT_sgAB_R2_sense",
      "WT_sgA_R1_branch", "WT_sgA_R1_cmv", "WT_sgA_R1_sense",
      "WT_sgB_R2_branch", "WT_sgB_R2_cmv", "WT_sgB_R2_sense",
    ],
    dtype = object,
  ).reshape((-1, 3))
  for i in range(image_grid.shape[0]):
    for j in range(image_grid.shape[1]):
      image_grid[i, j] = 'plots/graphs/individual/' + image_grid[i, j]

  legend_list = [
    # {
    #   'type': 'freq_ratio',
    #   'treatment_1': 'sense',
    #   'treatment_2': 'branch',
    #   'color_bar_file': 'images/freq_ratio_sense_cmv.png',
    # }
    {
      'type': 'variation_type',
    }
  ]
  make_slide(
    prs,
    title = '',
    image_grid_list = [image_grid],
    label_grid_list = [label_grid],
    node_size_max_freq = 1,
    node_size_min_freq = 1e-5,
    node_size_max_px = 200,
    node_size_min_px = 10,
    legend_list = legend_list,
  )
  prs.save('hello.pptx')
  # make_pptx_2()
  pass

